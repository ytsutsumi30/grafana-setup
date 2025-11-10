#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生産管理システム 操作手順書 生成スクリプト
index.html と maintenance.html 配下の画面・機能の操作手順書をPPTX形式で作成
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """タイトルスライドを作成"""
    slide_layout = prs.slide_layouts[0]  # タイトルスライドレイアウト
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # フォント設定
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)

    return slide

def create_section_header_slide(prs, title):
    """セクションヘッダースライドを作成"""
    slide_layout = prs.slide_layouts[6]  # 空白レイアウト
    slide = prs.slides.add_slide(slide_layout)

    # 背景色を設定（グラデーション風）
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 110, 253)  # Bootstrap primary blue

    # タイトルテキストボックス
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = title

    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(40)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    paragraph.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, title, content_items):
    """コンテンツスライドを作成"""
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツレイアウト
    slide = prs.slides.add_slide(slide_layout)

    # タイトル設定
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # コンテンツ領域
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()

        p = text_frame.paragraphs[i]

        if isinstance(item, dict):
            # ネストされた項目
            p.text = item.get('text', '')
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('font_size', 16))
            if item.get('bold'):
                p.font.bold = True
        else:
            # 通常のテキスト
            p.text = item
            p.font.size = Pt(16)

    return slide

def create_manual():
    """操作手順書を作成"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. タイトルスライド
    create_title_slide(
        prs,
        "生産管理システム 操作手順書",
        "index.html / maintenance.html 画面・機能ガイド\n\nバージョン 2.1.0"
    )

    # 2. 目次
    create_content_slide(prs, "目次", [
        {"text": "1. システム概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "2. 出荷検品システム (index.html)", "level": 0, "font_size": 18, "bold": True},
        {"text": "3. メンテナンス機能 (maintenance.html)", "level": 0, "font_size": 18, "bold": True},
        {"text": "4. マスタデータ管理", "level": 0, "font_size": 18, "bold": True},
        {"text": "5. 業務データ管理", "level": 0, "font_size": 18, "bold": True},
        {"text": "6. モニタリング・分析機能", "level": 0, "font_size": 18, "bold": True},
        {"text": "7. 品質管理（QCツール）", "level": 0, "font_size": 18, "bold": True},
        {"text": "8. システム管理", "level": 0, "font_size": 18, "bold": True}
    ])

    # 3. システム概要
    create_section_header_slide(prs, "1. システム概要")

    create_content_slide(prs, "システム概要", [
        {"text": "システム名称", "level": 0, "font_size": 18, "bold": True},
        {"text": "日本の生産管理システム（出荷検品・在庫管理）", "level": 1, "font_size": 16},
        "",
        {"text": "主要機能", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 出荷指示の検品実行（QRコード対応）", "level": 1, "font_size": 16},
        {"text": "• リアルタイム在庫管理", "level": 1, "font_size": 16},
        {"text": "• 製品・拠点・構成部品のマスタデータ管理", "level": 1, "font_size": 16},
        {"text": "• 品質管理ツール（QC七つ道具・新QC七つ道具）", "level": 1, "font_size": 16},
        {"text": "• データベースバックアップ・復元", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "技術スタック", [
        {"text": "フロントエンド", "level": 0, "font_size": 18, "bold": True},
        {"text": "• Vanilla JavaScript (ES6+), Bootstrap 5, HTML5", "level": 1, "font_size": 16},
        "",
        {"text": "バックエンド", "level": 0, "font_size": 18, "bold": True},
        {"text": "• Node.js 18+ (Express)", "level": 1, "font_size": 16},
        "",
        {"text": "データベース", "level": 0, "font_size": 18, "bold": True},
        {"text": "• PostgreSQL 15", "level": 1, "font_size": 16},
        "",
        {"text": "その他", "level": 0, "font_size": 18, "bold": True},
        {"text": "• Docker Compose によるコンテナ管理", "level": 1, "font_size": 16},
        {"text": "• nginx リバースプロキシ", "level": 1, "font_size": 16},
        {"text": "• Grafana + Prometheus (オプショナル)", "level": 1, "font_size": 16}
    ])

    # 4. 出荷検品システム (index.html)
    create_section_header_slide(prs, "2. 出荷検品システム")

    create_content_slide(prs, "出荷検品システム - 画面概要", [
        {"text": "アクセス方法", "level": 0, "font_size": 18, "bold": True},
        {"text": "URL: https://[ホスト名]/index.html", "level": 1, "font_size": 16},
        "",
        {"text": "主要機能エリア", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. ヘッダー部", "level": 1, "font_size": 16},
        {"text": "   - システムタイトル・コード修正日時表示", "level": 2, "font_size": 14},
        {"text": "   - QRコード読み取り・メンテナンス・メインシステムへのリンク", "level": 2, "font_size": 14},
        {"text": "2. 検品待ち出荷指示一覧", "level": 1, "font_size": 16},
        {"text": "   - ステータスが「pending」の出荷指示を表示", "level": 2, "font_size": 14},
        {"text": "3. サイドバー", "level": 1, "font_size": 16},
        {"text": "   - 本日の検品実績（完了・待機中・合格率）", "level": 2, "font_size": 14},
        {"text": "   - 注意事項・最近の検品履歴", "level": 2, "font_size": 14}
    ])

    create_content_slide(prs, "検品待ち出荷指示の操作", [
        {"text": "出荷指示カードの表示内容", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 出荷指示番号・製品名・製品コード", "level": 1, "font_size": 16},
        {"text": "• 出荷数量・顧客名・出荷日", "level": 1, "font_size": 16},
        {"text": "• 製品QRコード（クリックでQR検品画面へ遷移）", "level": 1, "font_size": 16},
        {"text": "• 配送先・特記事項", "level": 1, "font_size": 16},
        {"text": "• 優先度バッジ（高優先度・通常優先度・低優先度）", "level": 1, "font_size": 16},
        "",
        {"text": "操作ボタン", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 検品開始：従来型の検品フォームを開く", "level": 1, "font_size": 16},
        {"text": "• QR検品：QRコードベースの検品画面を開く", "level": 1, "font_size": 16},
        {"text": "• 詳細表示：出荷指示の詳細情報を表示", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "従来型検品の操作手順", [
        {"text": "Step 1: 検品開始", "level": 0, "font_size": 18, "bold": True},
        {"text": "出荷指示カードの「検品開始」ボタンをクリック", "level": 1, "font_size": 16},
        "",
        {"text": "Step 2: 検品情報の入力", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 検品者名を入力（必須）", "level": 1, "font_size": 16},
        {"text": "• 検品数量を入力（デフォルト: 予定数量）", "level": 1, "font_size": 16},
        {"text": "• 合格数量を入力（デフォルト: 予定数量）", "level": 1, "font_size": 16},
        {"text": "• 不合格数量は自動計算される", "level": 1, "font_size": 16},
        "",
        {"text": "Step 3: チェック項目の確認", "level": 0, "font_size": 18, "bold": True},
        {"text": "以下の項目をチェック:", "level": 1, "font_size": 16},
        {"text": "□ ラベル確認完了 □ 梱包状態確認 □ 出荷書類確認", "level": 1, "font_size": 14},
        {"text": "□ 品質基準適合 □ 数量一致確認 □ 最終承認", "level": 1, "font_size": 14}
    ])

    create_content_slide(prs, "従来型検品の操作手順（続き）", [
        {"text": "Step 4: 検品完了・下書き保存", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 下書き保存：検品途中で一時保存", "level": 1, "font_size": 16},
        {"text": "• 検品完了：検品を完了し、データベースに記録", "level": 1, "font_size": 16},
        {"text": "  - 最終承認が未チェックの場合、確認ダイアログが表示される", "level": 2, "font_size": 14},
        "",
        {"text": "注意事項", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 検品者名は必須項目です", "level": 1, "font_size": 16},
        {"text": "• 合格数量は検品数量を超えることはできません", "level": 1, "font_size": 16},
        {"text": "• 検品完了後、出荷指示のステータスが更新されます", "level": 1, "font_size": 16},
        {"text": "• 在庫は自動的に減算されます", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "QR検品の操作手順", [
        {"text": "Step 1: QR検品画面を開く", "level": 0, "font_size": 18, "bold": True},
        {"text": "出荷指示カードの「QR検品」ボタンまたは製品QRコードをクリック", "level": 1, "font_size": 16},
        {"text": "→ 別タブでQR検品画面（qr-inspection.html）が開く", "level": 1, "font_size": 16},
        "",
        {"text": "Step 2: 検品者名を入力", "level": 0, "font_size": 18, "bold": True},
        {"text": "QR検品画面左側の「検品者名」欄に名前を入力", "level": 1, "font_size": 16},
        "",
        {"text": "Step 3: QRスキャン開始", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 「QRスキャン開始」ボタンをクリック", "level": 1, "font_size": 16},
        {"text": "• カメラアクセスの許可を求められたら「許可」をクリック", "level": 1, "font_size": 16},
        {"text": "• カメラ映像が表示される", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "QR検品の操作手順（続き）", [
        {"text": "Step 4: 同梱物のQRコードをスキャン", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 製品構成部品のQRコードをカメラに向ける", "level": 1, "font_size": 16},
        {"text": "• スキャン成功すると、該当アイテムが「確認済み」に変わる", "level": 1, "font_size": 16},
        {"text": "• 進捗表示「N/M」が更新される", "level": 1, "font_size": 16},
        {"text": "• 次のアイテムのスキャンに自動的に進む", "level": 1, "font_size": 16},
        "",
        {"text": "Step 5: 検品完了", "level": 0, "font_size": 18, "bold": True},
        {"text": "• すべてのアイテムをスキャン後、「検品完了」ボタンをクリック", "level": 1, "font_size": 16},
        {"text": "• 検品結果がデータベースに記録される", "level": 1, "font_size": 16},
        {"text": "• 在庫が自動的に減算される", "level": 1, "font_size": 16},
        {"text": "• 完了メッセージが表示される", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "QR検品の補助機能", [
        {"text": "テストスキャン", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 「テストスキャン」ボタンでランダムにアイテムをスキャンシミュレート", "level": 1, "font_size": 16},
        {"text": "• 開発・デモ用途に使用", "level": 1, "font_size": 16},
        "",
        {"text": "手動入力", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 「手動入力」ボタンでQRコード値を手入力可能", "level": 1, "font_size": 16},
        {"text": "• カメラが使えない場合のフォールバック", "level": 1, "font_size": 16},
        "",
        {"text": "カメラの注意事項", "level": 0, "font_size": 18, "bold": True},
        {"text": "• HTTPSまたはlocalhostでアクセスすること", "level": 1, "font_size": 16},
        {"text": "• iOS Safariの場合、設定→Safari→カメラを「許可」に設定", "level": 1, "font_size": 16},
        {"text": "• ブラウザを最新版に更新すること", "level": 1, "font_size": 16}
    ])

    # 5. メンテナンス機能 (maintenance.html)
    create_section_header_slide(prs, "3. メンテナンス機能")

    create_content_slide(prs, "メンテナンス画面 - 概要", [
        {"text": "アクセス方法", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 出荷検品システム画面右上の「メンテナンス」ボタンをクリック", "level": 1, "font_size": 16},
        {"text": "• 直接アクセス: https://[ホスト名]/maintenance.html", "level": 1, "font_size": 16},
        "",
        {"text": "メンテナンス画面の構成", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. マスタデータ管理", "level": 1, "font_size": 16},
        {"text": "   - 製品マスタ・出荷元拠点・配送先拠点・製品構成部品", "level": 2, "font_size": 14},
        {"text": "2. 業務データ管理", "level": 1, "font_size": 16},
        {"text": "   - 生産計画・出荷指示・在庫管理", "level": 2, "font_size": 14},
        {"text": "3. モニタリング・分析", "level": 1, "font_size": 16},
        {"text": "   - モニタリングダッシュボード", "level": 2, "font_size": 14},
        {"text": "4. 品質管理（QCツール）", "level": 1, "font_size": 16},
        {"text": "   - QC七つ道具・新QC七つ道具", "level": 2, "font_size": 14}
    ])

    create_content_slide(prs, "メンテナンス画面 - 概要（続き）", [
        {"text": "メンテナンス画面の構成（続き）", "level": 0, "font_size": 18, "bold": True},
        {"text": "5. システム管理", "level": 1, "font_size": 16},
        {"text": "   - データベース・システムログ・ユーザー管理", "level": 2, "font_size": 14},
        "",
        {"text": "画面下部のシステム情報", "level": 0, "font_size": 18, "bold": True},
        {"text": "• システムバージョン: v2.1.0", "level": 1, "font_size": 16},
        {"text": "• データベース: PostgreSQL 15", "level": 1, "font_size": 16},
        {"text": "• テーブル数: 12", "level": 1, "font_size": 16},
        {"text": "• 最終更新: 2025-11-07", "level": 1, "font_size": 16}
    ])

    # 6. マスタデータ管理
    create_section_header_slide(prs, "4. マスタデータ管理")

    create_content_slide(prs, "製品マスタ (products.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "製品情報の登録・編集・削除を行います", "level": 1, "font_size": 16},
        "",
        {"text": "管理項目", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 製品コード（一意の識別子）", "level": 1, "font_size": 16},
        {"text": "• 製品名称", "level": 1, "font_size": 16},
        {"text": "• 単価", "level": 1, "font_size": 16},
        {"text": "• 在庫数量", "level": 1, "font_size": 16},
        {"text": "• 製品説明・備考", "level": 1, "font_size": 16},
        "",
        {"text": "操作方法", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 新規登録：「新規製品登録」ボタンをクリック", "level": 1, "font_size": 16},
        {"text": "2. 編集：製品行の「編集」ボタンをクリック", "level": 1, "font_size": 16},
        {"text": "3. 削除：製品行の「削除」ボタンをクリック", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "出荷元拠点 (shipping-locations.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "出荷元となる倉庫・工場の情報を管理します", "level": 1, "font_size": 16},
        "",
        {"text": "管理項目", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 拠点コード（一意の識別子）", "level": 1, "font_size": 16},
        {"text": "• 拠点名称", "level": 1, "font_size": 16},
        {"text": "• 住所（郵便番号・都道府県・市区町村・番地）", "level": 1, "font_size": 16},
        {"text": "• 連絡先（電話番号・担当者名）", "level": 1, "font_size": 16},
        "",
        {"text": "用途", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 出荷指示作成時に出荷元を選択", "level": 1, "font_size": 16},
        {"text": "• 出荷書類への印字", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "配送先拠点 (delivery-locations.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "配送先の顧客拠点情報を管理します", "level": 1, "font_size": 16},
        "",
        {"text": "管理項目", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 拠点コード（一意の識別子）", "level": 1, "font_size": 16},
        {"text": "• 拠点名称", "level": 1, "font_size": 16},
        {"text": "• 配送方法（トラック・航空便・船便など）", "level": 1, "font_size": 16},
        {"text": "• 住所情報", "level": 1, "font_size": 16},
        {"text": "• 配送業者情報", "level": 1, "font_size": 16},
        "",
        {"text": "用途", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 出荷指示作成時に配送先を選択", "level": 1, "font_size": 16},
        {"text": "• 配送ラベルへの印字", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "製品構成部品 (product-components.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "製品を構成する部品とQRコードのマッピングを管理します", "level": 1, "font_size": 16},
        "",
        {"text": "管理項目", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 製品ID（対象製品）", "level": 1, "font_size": 16},
        {"text": "• 部品種別（例: 本体、電源ケーブル、説明書、保証書）", "level": 1, "font_size": 16},
        {"text": "• 部品名称", "level": 1, "font_size": 16},
        {"text": "• QRコード値", "level": 1, "font_size": 16},
        {"text": "• 必須フラグ（検品時に必須かどうか）", "level": 1, "font_size": 16},
        "",
        {"text": "用途", "level": 0, "font_size": 18, "bold": True},
        {"text": "• QR検品時の同梱物チェックリスト生成", "level": 1, "font_size": 16},
        {"text": "• スキャンされたQRコードとの照合", "level": 1, "font_size": 16}
    ])

    # 7. 業務データ管理
    create_section_header_slide(prs, "5. 業務データ管理")

    create_content_slide(prs, "生産計画 (production-plans.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "製品の生産計画を管理します", "level": 1, "font_size": 16},
        "",
        {"text": "管理項目", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 計画番号", "level": 1, "font_size": 16},
        {"text": "• 製品ID（対象製品）", "level": 1, "font_size": 16},
        {"text": "• 計画数量", "level": 1, "font_size": 16},
        {"text": "• 開始予定日・完了予定日", "level": 1, "font_size": 16},
        {"text": "• 実績数量", "level": 1, "font_size": 16},
        {"text": "• ステータス（計画中・生産中・完了）", "level": 1, "font_size": 16},
        "",
        {"text": "操作方法", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 新規計画：製品と数量を入力", "level": 1, "font_size": 16},
        {"text": "2. 実績入力：生産完了数を更新", "level": 1, "font_size": 16},
        {"text": "3. 在庫連携：生産完了時に在庫が自動増加", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "出荷指示 (shipping-instructions.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "顧客への出荷指示を作成・管理します", "level": 1, "font_size": 16},
        "",
        {"text": "管理項目", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 出荷指示番号", "level": 1, "font_size": 16},
        {"text": "• 製品ID・出荷数量", "level": 1, "font_size": 16},
        {"text": "• 出荷元拠点・配送先拠点", "level": 1, "font_size": 16},
        {"text": "• 出荷日・配送方法", "level": 1, "font_size": 16},
        {"text": "• ステータス（pending / in_progress / completed / cancelled）", "level": 1, "font_size": 16},
        {"text": "• 優先度（high / normal / low）", "level": 1, "font_size": 16},
        {"text": "• 特記事項", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "出荷指示 - 操作フロー", [
        {"text": "新規出荷指示の作成", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 「新規出荷指示」ボタンをクリック", "level": 1, "font_size": 16},
        {"text": "2. 製品・数量・出荷日を入力", "level": 1, "font_size": 16},
        {"text": "3. 出荷元・配送先を選択", "level": 1, "font_size": 16},
        {"text": "4. 優先度を設定（必要に応じて）", "level": 1, "font_size": 16},
        {"text": "5. 「登録」ボタンで保存", "level": 1, "font_size": 16},
        "",
        {"text": "検品との連携", "level": 0, "font_size": 18, "bold": True},
        {"text": "• ステータスが「pending」の出荷指示がindex.htmlに表示される", "level": 1, "font_size": 16},
        {"text": "• 検品完了後、ステータスが「completed」に自動更新される", "level": 1, "font_size": 16},
        {"text": "• 在庫が自動的に減算される", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "在庫管理 (inventory.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "製品在庫の確認と調整を行います", "level": 1, "font_size": 16},
        "",
        {"text": "表示項目", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 製品コード・製品名", "level": 1, "font_size": 16},
        {"text": "• 現在庫数（current_stock）", "level": 1, "font_size": 16},
        {"text": "• 引当済数（reserved_stock）", "level": 1, "font_size": 16},
        {"text": "• 利用可能在庫（available_stock = 現在庫 - 引当済）", "level": 1, "font_size": 16},
        {"text": "  ※ available_stockはPostgreSQLのGENERATED列", "level": 1, "font_size": 14},
        "",
        {"text": "在庫調整", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 「在庫調整」ボタンで現在庫を手動で増減", "level": 1, "font_size": 16},
        {"text": "• 調整理由を入力（棚卸、不良品廃棄など）", "level": 1, "font_size": 16}
    ])

    # 8. モニタリング・分析機能
    create_section_header_slide(prs, "6. モニタリング・分析機能")

    create_content_slide(prs, "モニタリングダッシュボード (monitoring.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "リアルタイムで出荷・在庫・検品状況を監視します", "level": 1, "font_size": 16},
        "",
        {"text": "表示内容", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 出荷指示ステータス別集計（pending / completed / cancelled）", "level": 1, "font_size": 16},
        {"text": "• 在庫アラート（在庫切れ・低在庫警告）", "level": 1, "font_size": 16},
        {"text": "• 検品パフォーマンス（完了件数・合格率・平均処理時間）", "level": 1, "font_size": 16},
        {"text": "• 製品別出荷数量ランキング", "level": 1, "font_size": 16},
        {"text": "• 時系列グラフ（日次・週次・月次）", "level": 1, "font_size": 16},
        "",
        {"text": "用途", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 生産・出荷状況の可視化", "level": 1, "font_size": 16},
        {"text": "• ボトルネックの早期発見", "level": 1, "font_size": 16}
    ])

    # 9. 品質管理（QCツール）
    create_section_header_slide(prs, "7. 品質管理（QCツール）")

    create_content_slide(prs, "QC七つ道具 (qc-dashboard.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "数値データを用いた品質管理・統計分析ツール", "level": 1, "font_size": 16},
        "",
        {"text": "提供される手法", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. パレート図", "level": 1, "font_size": 16},
        {"text": "   - 不良原因の重点分析（80-20の法則）", "level": 2, "font_size": 14},
        {"text": "2. ヒストグラム", "level": 1, "font_size": 16},
        {"text": "   - データの分布・ばらつきの可視化", "level": 2, "font_size": 14},
        {"text": "3. 管理図（コントロールチャート）", "level": 1, "font_size": 16},
        {"text": "   - プロセスの安定性監視", "level": 2, "font_size": 14},
        {"text": "4. 散布図", "level": 1, "font_size": 16},
        {"text": "   - 2つの変数の相関関係分析", "level": 2, "font_size": 14},
        {"text": "5. チェックシート・特性要因図・層別分析", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "QC七つ道具 - 使用例", [
        {"text": "パレート図の作成手順", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 「パレート図」タブを選択", "level": 1, "font_size": 16},
        {"text": "2. 分析対象データ（不良原因・発生件数）を入力", "level": 1, "font_size": 16},
        {"text": "3. 「グラフ生成」ボタンをクリック", "level": 1, "font_size": 16},
        {"text": "4. パレート図が表示され、上位80%の原因が強調表示される", "level": 1, "font_size": 16},
        "",
        {"text": "管理図の作成手順", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 「管理図」タブを選択", "level": 1, "font_size": 16},
        {"text": "2. 時系列測定データを入力", "level": 1, "font_size": 16},
        {"text": "3. 管理限界線（UCL・LCL）を設定", "level": 1, "font_size": 16},
        {"text": "4. 管理図が表示され、異常値が検出される", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "新QC七つ道具 (qc-analysis.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "定性データを用いた課題整理・問題解決ツール", "level": 1, "font_size": 16},
        "",
        {"text": "提供される手法", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 親和図法（KJ法）", "level": 1, "font_size": 16},
        {"text": "   - 言語データをグループ化して問題の全体像を把握", "level": 2, "font_size": 14},
        {"text": "2. 連関図法", "level": 1, "font_size": 16},
        {"text": "   - 複雑な要因の因果関係を可視化", "level": 2, "font_size": 14},
        {"text": "3. 系統図法", "level": 1, "font_size": 16},
        {"text": "   - 目的達成のための手段を階層的に整理", "level": 2, "font_size": 14},
        {"text": "4. マトリックス図法・アローダイアグラム", "level": 1, "font_size": 16},
        {"text": "5. PDPC法（プロセス決定計画図）", "level": 1, "font_size": 16},
        {"text": "6. マトリックスデータ解析法", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "新QC七つ道具 - 使用例", [
        {"text": "親和図法の作成手順", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 「親和図法」タブを選択", "level": 1, "font_size": 16},
        {"text": "2. 問題に関する意見・データをカードに入力", "level": 1, "font_size": 16},
        {"text": "3. 類似したカードをドラッグ&ドロップでグループ化", "level": 1, "font_size": 16},
        {"text": "4. グループごとに見出しを付ける", "level": 1, "font_size": 16},
        {"text": "5. 親和図が完成し、問題の全体像が可視化される", "level": 1, "font_size": 16},
        "",
        {"text": "系統図法の作成手順", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 「系統図法」タブを選択", "level": 1, "font_size": 16},
        {"text": "2. 目的（最終目標）を入力", "level": 1, "font_size": 16},
        {"text": "3. 手段を階層的に追加していく", "level": 1, "font_size": 16},
        {"text": "4. 系統図が完成し、実行計画が明確になる", "level": 1, "font_size": 16}
    ])

    # 10. システム管理
    create_section_header_slide(prs, "8. システム管理")

    create_content_slide(prs, "データベース管理 (database.html)", [
        {"text": "機能概要", "level": 0, "font_size": 18, "bold": True},
        {"text": "PostgreSQLデータベースのバックアップ・復元を行います", "level": 1, "font_size": 16},
        "",
        {"text": "バックアップ", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 「バックアップ実行」ボタンをクリック", "level": 1, "font_size": 16},
        {"text": "2. pg_dumpコマンドでデータベース全体をダンプ", "level": 1, "font_size": 16},
        {"text": "3. バックアップファイルがbackups/ディレクトリに保存される", "level": 1, "font_size": 16},
        {"text": "4. ファイル名形式: backup_YYYYMMDD_HHMMSS.sql", "level": 1, "font_size": 16},
        "",
        {"text": "復元", "level": 0, "font_size": 18, "bold": True},
        {"text": "1. 「復元」セクションでバックアップファイルを選択", "level": 1, "font_size": 16},
        {"text": "2. 「復元実行」ボタンをクリック", "level": 1, "font_size": 16},
        {"text": "3. 警告メッセージを確認後、実行を承認", "level": 1, "font_size": 16},
        {"text": "4. psqlコマンドでデータベースが復元される", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "データベース管理 - 注意事項", [
        {"text": "バックアップの推奨頻度", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 日次バックアップ（営業日終了時）", "level": 1, "font_size": 16},
        {"text": "• 重要な操作前（マスタデータ大量更新など）", "level": 1, "font_size": 16},
        {"text": "• システムアップデート前", "level": 1, "font_size": 16},
        "",
        {"text": "復元時の注意事項", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 復元は既存データを上書きします", "level": 1, "font_size": 16},
        {"text": "• 復元前に必ず現在のデータをバックアップしてください", "level": 1, "font_size": 16},
        {"text": "• 復元中はシステムを使用しないでください", "level": 1, "font_size": 16},
        {"text": "• 復元後、データ整合性を確認してください", "level": 1, "font_size": 16},
        "",
        {"text": "コマンドラインからのバックアップ", "level": 0, "font_size": 18, "bold": True},
        {"text": "./manage.sh backup", "level": 1, "font_size": 14}
    ])

    create_content_slide(prs, "システムログ・ユーザー管理", [
        {"text": "システムログ", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 現在準備中", "level": 1, "font_size": 16},
        {"text": "• 将来的にアプリケーションログ・エラーログの閲覧が可能になります", "level": 1, "font_size": 16},
        "",
        {"text": "ユーザー管理", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 現在準備中", "level": 1, "font_size": 16},
        {"text": "• 将来的にユーザーアカウント・権限管理が可能になります", "level": 1, "font_size": 16},
        {"text": "• 予定機能:", "level": 1, "font_size": 16},
        {"text": "  - ユーザー登録・編集・削除", "level": 2, "font_size": 14},
        {"text": "  - ロール管理（管理者・検品者・閲覧者など）", "level": 2, "font_size": 14},
        {"text": "  - アクセス制御・ログイン履歴", "level": 2, "font_size": 14}
    ])

    # 11. 付録・トラブルシューティング
    create_section_header_slide(prs, "付録・トラブルシューティング")

    create_content_slide(prs, "よくある質問（FAQ）", [
        {"text": "Q1: QRカメラが起動しません", "level": 0, "font_size": 16, "bold": True},
        {"text": "A: HTTPSまたはlocalhostでアクセスしてください。", "level": 1, "font_size": 14},
        {"text": "   iOS Safariの場合、設定→Safari→カメラを「許可」に設定。", "level": 1, "font_size": 14},
        "",
        {"text": "Q2: 検品完了後、在庫が減らない", "level": 0, "font_size": 16, "bold": True},
        {"text": "A: QR検品の場合、「検品完了」ボタンを押す必要があります。", "level": 1, "font_size": 14},
        {"text": "   従来型検品の場合、「検品完了」ボタンで在庫が自動減算されます。", "level": 1, "font_size": 14},
        "",
        {"text": "Q3: 出荷指示が検品待ち一覧に表示されない", "level": 0, "font_size": 16, "bold": True},
        {"text": "A: ステータスが「pending」になっているか確認してください。", "level": 1, "font_size": 14},
        {"text": "   maintenance.html → 出荷指示 で確認・編集できます。", "level": 1, "font_size": 14},
        "",
        {"text": "Q4: 製品構成部品のQRコードを変更したい", "level": 0, "font_size": 16, "bold": True},
        {"text": "A: maintenance.html → 製品構成部品 で編集できます。", "level": 1, "font_size": 14}
    ])

    create_content_slide(prs, "システム起動・停止コマンド", [
        {"text": "システム起動", "level": 0, "font_size": 18, "bold": True},
        {"text": "./manage.sh start", "level": 1, "font_size": 16},
        {"text": "nginx + API + PostgreSQL が起動します", "level": 1, "font_size": 14},
        "",
        {"text": "システム停止", "level": 0, "font_size": 18, "bold": True},
        {"text": "./manage.sh stop", "level": 1, "font_size": 16},
        "",
        {"text": "システム再起動", "level": 0, "font_size": 18, "bold": True},
        {"text": "./manage.sh restart", "level": 1, "font_size": 16},
        "",
        {"text": "モニタリング起動（Grafana + Prometheus）", "level": 0, "font_size": 18, "bold": True},
        {"text": "./manage.sh monitoring start", "level": 1, "font_size": 16},
        "",
        {"text": "ログ確認", "level": 0, "font_size": 18, "bold": True},
        {"text": "./manage.sh logs", "level": 1, "font_size": 16}
    ])

    create_content_slide(prs, "お問い合わせ・サポート", [
        {"text": "技術サポート", "level": 0, "font_size": 18, "bold": True},
        {"text": "システムに関するご質問・トラブルシューティングは", "level": 1, "font_size": 16},
        {"text": "システム管理者までお問い合わせください。", "level": 1, "font_size": 16},
        "",
        {"text": "システム情報", "level": 0, "font_size": 18, "bold": True},
        {"text": "• バージョン: v2.1.0", "level": 1, "font_size": 16},
        {"text": "• データベース: PostgreSQL 15", "level": 1, "font_size": 16},
        {"text": "• 最終更新: 2025-11-07", "level": 1, "font_size": 16},
        "",
        {"text": "ドキュメント", "level": 0, "font_size": 18, "bold": True},
        {"text": "詳細な技術情報は CLAUDE.md を参照してください。", "level": 1, "font_size": 16},
        "",
        {"text": "本操作手順書の更新履歴", "level": 0, "font_size": 18, "bold": True},
        {"text": "• 2025-11-10: 初版作成", "level": 1, "font_size": 16}
    ])

    # ファイル保存
    output_path = "/home/user/grafana-setup/操作手順書_生産管理システム.pptx"
    prs.save(output_path)
    print(f"✓ PPTXファイルを作成しました: {output_path}")
    return output_path

if __name__ == "__main__":
    create_manual()
